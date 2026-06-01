#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
将用户提供的内容填入Y038V划伤LOP报告范本PPT - 修复版
使用python-pptx直接操作，保留范本格式
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
import copy
import os

# ── 路径 ──
ORIGINAL = r'D:\Desktop\Y038V客户2026年05月30日铜箔（划伤）问题LOP报告.pptx'
OUTPUT   = r'D:\Desktop\Y038V客户2026年05月30日铜箔（划伤）问题LOP报告.pptx'

prs = Presentation(ORIGINAL)

def clear_cell(cell):
    """清空单元格所有段落，保留第一个段落"""
    for i in range(len(cell.text_frame.paragraphs) - 1, 0, -1):
        p = cell.text_frame.paragraphs[i]._p
        p.getparent().remove(p)
    p = cell.text_frame.paragraphs[0]
    for r in list(p.runs):
        p._p.remove(r._r)

def set_cell_text(cell, text, font_size=Pt(10), bold=None):
    """设置单元格文本"""
    clear_cell(cell)
    p = cell.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    if font_size:
        run.font.size = font_size
    if bold is not None:
        run.font.bold = bold
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def set_cell_multiline(cell, lines, font_size=Pt(10), bold_first=False):
    """设置单元格多行文本"""
    clear_cell(cell)
    tf = cell.text_frame
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        if font_size:
            run.font.size = font_size
        if bold_first and i == 0:
            run.font.bold = True
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

def add_row_to_table(table, cell_texts, font_size=Pt(10)):
    """在表格末尾添加一行并填入文本"""
    rows = list(table.rows)
    if not rows:
        return
    new_row = copy.deepcopy(rows[-1]._tr)
    table._tbl.append(new_row)
    
    row_idx = len(list(table.rows)) - 1
    for col_idx, text in enumerate(cell_texts):
        if col_idx < len(table.columns):
            cell = table.cell(row_idx, col_idx)
            set_cell_text(cell, text, font_size=font_size)

def clear_textbox(shape):
    """清空文本框"""
    tf = shape.text_frame
    for i in range(len(tf.paragraphs) - 1, 0, -1):
        p = tf.paragraphs[i]._p
        p.getparent().remove(p)
    p = tf.paragraphs[0]
    for r in list(p.runs):
        p._p.remove(r._r)

def get_table_by_name(slide, name):
    """找到slide中指定名称的表格"""
    for shape in slide.shapes:
        if shape.has_table and shape.name == name:
            return shape.table
    return None

# ═══════════════════════════════════════════════════════════
# Slide 1: 问题描述 + 原因分析 + 临时围堵 + 改善措施
# ═══════════════════════════════════════════════════════════
slide1 = prs.slides[0]
table11 = None
for shape in slide1.shapes:
    if shape.name == "表格 11":
        table11 = shape.table
        break

if table11:
    # Cell[0,0] = 问题描述
    cell = table11.cell(0, 0)
    clear_cell(cell)
    tf = cell.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "问题描述"
    run.font.bold = True
    run.font.size = Pt(14)
    
    items = [
        ("When：", "2026年05月30日"),
        ("Where：", "时代长安来料检验车间"),
        ("Who：", "时代长安来料检验人员"),
        ("What：", "双光6*1390mm-ZQ铜箔，卷号A-264109221，批次0024328554，客户拆包检验发现铜箔端面明显划伤"),
        ("Why：", "分切、收卷、检验均无异常，仅包装后出现；端面受手指/手套摩擦导致"),
        ("How：", "影响客户上料与产线效率"),
        ("How much：", "1卷不良；非重复性问题"),
    ]
    for label, value in items:
        p = tf.add_paragraph()
        run_label = p.add_run()
        run_label.text = label
        run_label.font.bold = True
        run_label.font.size = Pt(10)
        run_val = p.add_run()
        run_val.text = value
        run_val.font.size = Pt(10)
    
    # Cell[0,1] = 原因分析
    cell = table11.cell(0, 1)
    clear_cell(cell)
    tf = cell.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "原因分析"
    run.font.bold = True
    run.font.size = Pt(14)
    
    reasons = [
        ("1. 发生原因", "包装工位人员徒手/佩戴不洁/破损手套，在扶正、对齐、套护角时手指/指甲直接触碰并摩擦铜箔端面，形成不规则划伤。"),
        ("2. 流出原因", "出厂终检未对端面100%全检，仅抽检表面，人为划伤漏检流出。"),
        ("3. 系统原因", "包装SOP未明确禁止触碰端面、手套标准、操作姿势；无物理防护措施；培训与监督不到位。"),
    ]
    for label, value in reasons:
        p = tf.add_paragraph()
        run_label = p.add_run()
        run_label.text = label
        run_label.font.bold = True
        run_label.font.size = Pt(10)
        p2 = tf.add_paragraph()
        run_val = p2.add_run()
        run_val.text = value
        run_val.font.size = Pt(9)
    
    # Cell[1,0] = 临时/围堵
    cell = table11.cell(1, 0)
    clear_cell(cell)
    tf = cell.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "临时/围堵"
    run.font.bold = True
    run.font.size = Pt(14)
    
    containment = [
        ("1. 风险范围锁定", "锁定同批次0024328554 + 同包装班组 + 前后3天出货所有卷料。"),
        ("2. 物料围堵", ["客户线边：全检、隔离、退换", "内/外仓：100%全检端面，不合格隔离", "在途：到货后全检放行", "供应商仓：全检合格再出货", "半成品：暂停出货，全检合格后放行"]),
        ("3. 恢复生产临时措施", ["包装岗立即停线整改，全员佩戴洁净丁腈手套", "严禁徒手/指甲触碰铜箔端面", "出货前增加端面100%全检，合格盖章放行"]),
    ]
    for label, value in containment:
        p = tf.add_paragraph()
        run_label = p.add_run()
        run_label.text = label
        run_label.font.bold = True
        run_label.font.size = Pt(10)
        if isinstance(value, list):
            for line in value:
                p2 = tf.add_paragraph()
                run_val = p2.add_run()
                run_val.text = line
                run_val.font.size = Pt(9)
        else:
            p2 = tf.add_paragraph()
            run_val = p2.add_run()
            run_val.text = value
            run_val.font.size = Pt(9)
    
    # Cell[1,1] = 改善措施
    cell = table11.cell(1, 1)
    clear_cell(cell)
    tf = cell.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "改善措施"
    run.font.bold = True
    run.font.size = Pt(14)
    
    improvements = [
        ("1. 发生改善", ["包装只许托卷芯/筒身，严禁手指接触端面", "统一使用洁净无粉丁腈手套，破损立即更换", "端面增加软质保护垫/无纺布隔离", "工位定时清铜粉、除异物"]),
        ("2. 流出改善", ["终检增加端面全检项目，双人复核", "不合格禁止流入下工序", "建立不良追溯台账"]),
        ("3. 系统改善", ["修订包装SOP，明确防触碰规范", "包装/检验岗专项培训+考核", "品管现场巡检，违规立即纠正", "更新FMEA与控制计划"]),
        ("4. 断点批次号", "2026-05-31起生产/包装批次为改善断点"),
    ]
    for label, value in improvements:
        p = tf.add_paragraph()
        run_label = p.add_run()
        run_label.text = label
        run_label.font.bold = True
        run_label.font.size = Pt(10)
        if isinstance(value, list):
            for line in value:
                p2 = tf.add_paragraph()
                run_val = p2.add_run()
                run_val.text = line
                run_val.font.size = Pt(9)
        else:
            p2 = tf.add_paragraph()
            run_val = p2.add_run()
            run_val.text = value
            run_val.font.size = Pt(9)

print('Slide 1 done')

# ═══════════════════════════════════════════════════════════
# Slide 2: 发生机理分析
# ═══════════════════════════════════════════════════════════
slide2 = prs.slides[1]
left = Emu(911424)
top = Emu(1500000)
width = Emu(10000000)
height = Emu(3000000)
txBox = slide2.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
run = p.add_run()
run.text = "发生机理："
run.font.bold = True
run.font.size = Pt(16)

p = tf.add_paragraph()
run = p.add_run()
run.text = "包装人员徒手/不洁手套  →  手指/指甲/异物接触铜箔端面  →  按压+滑动摩擦  →  铜箔端面表层形成不规则划伤"
run.font.size = Pt(14)

print('Slide 2 done')

# ═══════════════════════════════════════════════════════════
# Slide 3: 工序排查 - 填入圆角矩形文本
# ═══════════════════════════════════════════════════════════
slide3 = prs.slides[2]
process_steps_top = ['分切', '收卷', '成品检验', '包装', '仓储', '物流']
process_steps_bottom = ['分切', '收卷', '成品检验', '包装', '仓储', '物流', '客户端']

# 找所有空白的圆角矩形，按位置分行
auto_shapes = []
for shape in slide3.shapes:
    if shape.shape_type == 1 and shape.has_text_frame and not shape.text.strip():
        auto_shapes.append(shape)

row1 = sorted([s for s in auto_shapes if s.top < Emu(2500000)], key=lambda s: s.left)
row2 = sorted([s for s in auto_shapes if Emu(2500000) <= s.top < Emu(4500000)], key=lambda s: s.left)
row3 = sorted([s for s in auto_shapes if s.top >= Emu(4500000)], key=lambda s: s.left)

for i, shape in enumerate(row1):
    if i < len(process_steps_top):
        for para in shape.text_frame.paragraphs:
            if not para.text.strip():
                run = para.add_run()
                run.text = process_steps_top[i]
                run.font.size = Pt(10)
                run.font.bold = True
                break

for i, shape in enumerate(row2):
    if i < len(process_steps_top):
        for para in shape.text_frame.paragraphs:
            if not para.text.strip():
                run = para.add_run()
                run.text = process_steps_top[i]
                run.font.size = Pt(10)
                run.font.bold = True
                break

for i, shape in enumerate(row3):
    if i < len(process_steps_bottom):
        for para in shape.text_frame.paragraphs:
            if not para.text.strip():
                run = para.add_run()
                run.text = process_steps_bottom[i]
                run.font.size = Pt(10)
                run.font.bold = True
                break

# 填写小结
for shape in slide3.shapes:
    if shape.has_text_frame:
        txt = shape.text.strip()
        if txt == "小结：" or txt == "小结":
            p = shape.text_frame.add_paragraph()
            run = p.add_run()
            run.text = "异常仅在包装环节产生，因终检漏检流向客户。"
            run.font.size = Pt(10)
        elif txt == "发生工站":
            p = shape.text_frame.add_paragraph()
            run = p.add_run()
            run.text = "成品包装工位"
            run.font.size = Pt(9)
            run.font.bold = True
        elif txt == "流出工站":
            p = shape.text_frame.add_paragraph()
            run = p.add_run()
            run.text = "出厂终检工位"
            run.font.size = Pt(9)
            run.font.bold = True

print('Slide 3 done')

# ═══════════════════════════════════════════════════════════
# Slide 4: 发生原因分析表
# ═══════════════════════════════════════════════════════════
slide4 = prs.slides[3]
table4 = get_table_by_name(slide4, "表格 4")

if table4:
    rows = list(table4.rows)
    if len(rows) > 1:
        row1_data = ["1", "分切刀具/导辊不良", "分切收卷后端面完好", "非真因", "收卷无异常"]
        for col_idx, text in enumerate(row1_data):
            cell = table4.cell(1, col_idx)
            set_cell_text(cell, text, font_size=Pt(10))
    
    add_row_to_table(table4, ["2", "包装材料摩擦", "痕迹形态与辅材摩擦不符", "非真因", "无连续划痕"])
    add_row_to_table(table4, ["3", "手指/手套触碰端面", "模拟操作可复现相同划痕", "是真因", "与不良痕迹完全一致"])

print('Slide 4 done')

# ═══════════════════════════════════════════════════════════
# Slide 5: 流出原因分析表
# ═══════════════════════════════════════════════════════════
slide5 = prs.slides[4]
table5 = get_table_by_name(slide5, "表格 4")

if table5:
    rows = list(table5.rows)
    if len(rows) > 1:
        row1_data = ["1", "设备检测异常", "无在线端面检测设备", "非真因", "人工检验为主"]
        for col_idx, text in enumerate(row1_data):
            cell = table5.cell(1, col_idx)
            set_cell_text(cell, text, font_size=Pt(10))
    
    add_row_to_table(table5, ["2", "终检未100%全检端面", "原流程只抽检表面，不查端面", "是真因", "人为划伤漏检"])

print('Slide 5 done')

# ═══════════════════════════════════════════════════════════
# Slide 6: 发生原因（真因）
# ═══════════════════════════════════════════════════════════
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if shape.has_text_frame and "发生原因分析" in shape.text:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "发生原因分析" in run.text:
                    run.text = "发生原因（真因）"

table6 = get_table_by_name(slide6, "表格 4")
if table6:
    rows = list(table6.rows)
    if len(rows) > 1:
        row_data = ["3", "手指/手套触碰端面", "模拟操作可复现相同划痕", "是真因", "与不良痕迹完全一致"]
        for col_idx, text in enumerate(row_data):
            cell = table6.cell(1, col_idx)
            set_cell_text(cell, text, font_size=Pt(10))

# 添加真因结论
txBox = slide6.shapes.add_textbox(Emu(911424), Emu(700000), Emu(10000000), Emu(300000))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "包装人员在打包时，手指/手套直接触碰并摩擦铜箔端面，造成端面划伤。"
run.font.size = Pt(14)
run.font.bold = True

print('Slide 6 done')

# ═══════════════════════════════════════════════════════════
# Slide 7: 流出原因（真因）
# ═══════════════════════════════════════════════════════════
slide7 = prs.slides[6]
for shape in slide7.shapes:
    if shape.has_text_frame and "流出原因分析" in shape.text:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "流出原因分析" in run.text:
                    run.text = "流出原因（真因）"

table7 = get_table_by_name(slide7, "表格 4")
if table7:
    rows = list(table7.rows)
    if len(rows) > 1:
        row_data = ["4", "终检未100%全检端面", "原流程只抽检表面，不查端面", "是真因", "人为划伤漏检"]
        for col_idx, text in enumerate(row_data):
            cell = table7.cell(1, col_idx)
            set_cell_text(cell, text, font_size=Pt(10))

txBox = slide7.shapes.add_textbox(Emu(911424), Emu(700000), Emu(10000000), Emu(300000))
tf = txBox.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
run = p.add_run()
run.text = "出厂终检未对端面全检，导致不良品漏检流出。"
run.font.size = Pt(14)
run.font.bold = True

print('Slide 7 done')

# ═══════════════════════════════════════════════════════════
# Slide 8: 发生原因改善
# ═══════════════════════════════════════════════════════════
slide8 = prs.slides[7]
table8_s8 = get_table_by_name(slide8, "表格 8")
if table8_s8:
    rows = list(table8_s8.rows)
    if len(rows) > 1:
        row_data = ["人为操作", "包装时手指/手套触碰端面",
                    "1.严禁触碰端面，仅托卷芯\n2.统一洁净丁腈手套\n3.端面加软保护垫隔离\n4.工位5S定时清铜粉",
                    "包装组长", "2026-05-31", "2026-05-31", "已完成"]
        for col_idx, text in enumerate(row_data):
            if col_idx < len(table8_s8.columns):
                cell = table8_s8.cell(1, col_idx)
                set_cell_text(cell, text, font_size=Pt(9))

table4_s8 = get_table_by_name(slide8, "表格 4")
if table4_s8:
    cell = table4_s8.cell(0, 0)
    set_cell_multiline(cell, ["发生措施效果确认", "发生原因改善：人为操作-包装时手指/手套触碰端面"], font_size=Pt(10), bold_first=True)
    
    set_cell_text(table4_s8.cell(1, 0), "改善前：可徒手接触端面，无防护", font_size=Pt(10))
    set_cell_text(table4_s8.cell(1, 1), "改善后：物理隔离+规范动作，杜绝触碰", font_size=Pt(10))

print('Slide 8 done')

# ═══════════════════════════════════════════════════════════
# Slide 9: 流出原因改善
# ═══════════════════════════════════════════════════════════
slide9 = prs.slides[8]
table8_s9 = get_table_by_name(slide9, "表格 8")
if table8_s9:
    rows = list(table8_s9.rows)
    if len(rows) > 1:
        row_data = ["检验漏检", "终检未全检端面",
                    "1.端面100%全检并盖章\n2.包装后双人复核\n3.不良禁止放行",
                    "QC主管", "2026-05-31", "2026-05-31", "已完成"]
        for col_idx, text in enumerate(row_data):
            if col_idx < len(table8_s9.columns):
                cell = table8_s9.cell(1, col_idx)
                set_cell_text(cell, text, font_size=Pt(9))

table4_s9 = get_table_by_name(slide9, "表格 4")
if table4_s9:
    cell = table4_s9.cell(0, 0)
    set_cell_multiline(cell, ["流出措施效果确认", "流出原因改善：检验漏检-终检未全检端面"], font_size=Pt(10), bold_first=True)
    
    set_cell_text(table4_s9.cell(1, 0), "改善前：仅抽检表面，不查端面", font_size=Pt(10))
    set_cell_text(table4_s9.cell(1, 1), "改善后：端面必检，漏检率0%", font_size=Pt(10))

print('Slide 9 done')

# ═══════════════════════════════════════════════════════════
# Slide 10: 系统原因改善
# ═══════════════════════════════════════════════════════════
slide10 = prs.slides[9]
table8_s10 = get_table_by_name(slide10, "表格 8")
if table8_s10:
    rows = list(table8_s10.rows)
    if len(rows) > 1:
        row_data = ["管理缺失", "SOP不完善、培训不足",
                    "1.修订包装SOP防划伤条款\n2.全员培训+上岗考核\n3.品管巡检监督\n4.更新FMEA/控制计划",
                    "品质经理", "2026-05-31", "2026-05-31", "已完成"]
        for col_idx, text in enumerate(row_data):
            if col_idx < len(table8_s10.columns):
                cell = table8_s10.cell(1, col_idx)
                set_cell_text(cell, text, font_size=Pt(9))

table4_s10 = get_table_by_name(slide10, "表格 4")
if table4_s10:
    cell = table4_s10.cell(0, 0)
    set_cell_multiline(cell, ["系统措施效果确认", "系统原因改善：管理缺失-SOP不完善、培训不足"], font_size=Pt(10), bold_first=True)
    
    set_cell_text(table4_s10.cell(1, 0), "改善前：无明确规范，依赖自觉", font_size=Pt(10))
    set_cell_text(table4_s10.cell(1, 1), "改善后：制度+培训+监督闭环", font_size=Pt(10))

print('Slide 10 done')

# ═══════════════════════════════════════════════════════════
# Slide 11: 发生改善效果验证
# ═══════════════════════════════════════════════════════════
slide11 = prs.slides[10]
table8_s11 = get_table_by_name(slide11, "表格 8")
if table8_s11:
    rows = list(table8_s11.rows)
    if len(rows) > 1:
        row_data = ["人为操作", "包装时手指/手套触碰端面",
                    "1.严禁触碰端面，仅托卷芯\n2.统一洁净丁腈手套\n3.端面加软保护垫隔离\n4.工位5S定时清铜粉",
                    "包装组长", "2026-05-31", "2026-05-31", "已完成"]
        for col_idx, text in enumerate(row_data):
            if col_idx < len(table8_s11.columns):
                cell = table8_s11.cell(1, col_idx)
                set_cell_text(cell, text, font_size=Pt(9))

table4_s11 = get_table_by_name(slide11, "表格 4")
if table4_s11:
    cell = table4_s11.cell(0, 0)
    set_cell_multiline(cell, ["发生措施效果确认", "发生原因改善：人为操作-包装时手指/手套触碰端面"], font_size=Pt(10), bold_first=True)
    
    rows11 = list(table4_s11.rows)
    if len(rows11) > 1:
        set_cell_text(table4_s11.cell(1, 0), "验证过程", font_size=Pt(10), bold=True)
        set_cell_text(table4_s11.cell(1, 1), "验证效果说明", font_size=Pt(10), bold=True)
        set_cell_text(table4_s11.cell(1, 2), "图示", font_size=Pt(10), bold=True)
    
    if len(rows11) > 2:
        set_cell_text(table4_s11.cell(2, 0), "连续3天包装全流程监控", font_size=Pt(9))
        set_cell_text(table4_s11.cell(2, 1), "手套合规率100%，无触碰端面行为，无新增划伤", font_size=Pt(9))
        set_cell_text(table4_s11.cell(2, 2), "附改善后现场照片", font_size=Pt(9))
    
    if len(rows11) > 3:
        set_cell_text(table4_s11.cell(3, 0), "模拟包装对比测试", font_size=Pt(9))
        set_cell_text(table4_s11.cell(3, 1), "改善后无手指划伤产生", font_size=Pt(9))
        set_cell_text(table4_s11.cell(3, 2), "附OK端面照片", font_size=Pt(9))

print('Slide 11 done')

# ═══════════════════════════════════════════════════════════
# Slide 12: 流出改善效果验证
# ═══════════════════════════════════════════════════════════
slide12 = prs.slides[11]
table8_s12 = get_table_by_name(slide12, "表格 8")
if table8_s12:
    rows = list(table8_s12.rows)
    if len(rows) > 1:
        row_data = ["检验漏检", "终检未全检端面",
                    "1.端面100%全检并盖章\n2.包装后双人复核\n3.不良禁止放行",
                    "QC主管", "2026-05-31", "2026-05-31", "已完成"]
        for col_idx, text in enumerate(row_data):
            if col_idx < len(table8_s12.columns):
                cell = table8_s12.cell(1, col_idx)
                set_cell_text(cell, text, font_size=Pt(9))

table4_s12 = get_table_by_name(slide12, "表格 4")
if table4_s12:
    cell = table4_s12.cell(0, 0)
    set_cell_multiline(cell, ["流出措施效果确认", "流出原因改善：检验漏检-终检未全检端面"], font_size=Pt(10), bold_first=True)
    
    rows12 = list(table4_s12.rows)
    if len(rows12) > 1:
        set_cell_text(table4_s12.cell(1, 0), "验证过程", font_size=Pt(10), bold=True)
        set_cell_text(table4_s12.cell(1, 1), "验证效果说明", font_size=Pt(10), bold=True)
        set_cell_text(table4_s12.cell(1, 2), "图示", font_size=Pt(10), bold=True)
    
    if len(rows12) > 2:
        set_cell_text(table4_s12.cell(2, 0), "终检100%全检端面执行", font_size=Pt(9))
        set_cell_text(table4_s12.cell(2, 1), "连续批次零漏检，不良拦截率100%", font_size=Pt(9))
        set_cell_text(table4_s12.cell(2, 2), "附检验记录照片", font_size=Pt(9))
    
    if len(rows12) > 3:
        set_cell_text(table4_s12.cell(3, 0), "客户端退货/不良统计", font_size=Pt(9))
        set_cell_text(table4_s12.cell(3, 1), "无同类问题再次发生", font_size=Pt(9))
        set_cell_text(table4_s12.cell(3, 2), "附客户反馈记录", font_size=Pt(9))

print('Slide 12 done')

# ═══════════════════════════════════════════════════════════
# Slide 13: 系统改善效果验证
# ═══════════════════════════════════════════════════════════
slide13 = prs.slides[12]
table8_s13 = get_table_by_name(slide13, "表格 8")
if table8_s13:
    rows = list(table8_s13.rows)
    if len(rows) > 1:
        row_data = ["管理缺失", "SOP不完善、培训不足",
                    "1.修订包装SOP防划伤条款\n2.全员培训+上岗考核\n3.品管巡检监督\n4.更新FMEA/控制计划",
                    "品质经理", "2026-05-31", "2026-05-31", "已完成"]
        for col_idx, text in enumerate(row_data):
            if col_idx < len(table8_s13.columns):
                cell = table8_s13.cell(1, col_idx)
                set_cell_text(cell, text, font_size=Pt(9))

table4_s13 = get_table_by_name(slide13, "表格 4")
if table4_s13:
    cell = table4_s13.cell(0, 0)
    set_cell_multiline(cell, ["系统措施效果确认", "系统原因改善：管理缺失-SOP不完善、培训不足"], font_size=Pt(10), bold_first=True)
    
    rows13 = list(table4_s13.rows)
    if len(rows13) > 1:
        set_cell_text(table4_s13.cell(1, 0), "验证过程", font_size=Pt(10), bold=True)
        set_cell_text(table4_s13.cell(1, 1), "验证效果说明", font_size=Pt(10), bold=True)
        set_cell_text(table4_s13.cell(1, 2), "图示", font_size=Pt(10), bold=True)
    
    if len(rows13) > 2:
        set_cell_text(table4_s13.cell(2, 0), "SOP发布、培训考核", font_size=Pt(9))
        set_cell_text(table4_s13.cell(2, 1), "全员通过考核，执行到位", font_size=Pt(9))
        set_cell_text(table4_s13.cell(2, 2), "附培训签到/试卷", font_size=Pt(9))
    
    if len(rows13) > 3:
        set_cell_text(table4_s13.cell(3, 0), "现场巡检合规率", font_size=Pt(9))
        set_cell_text(table4_s13.cell(3, 1), "操作合规率100%", font_size=Pt(9))
        set_cell_text(table4_s13.cell(3, 2), "附巡检记录", font_size=Pt(9))

print('Slide 13 done')

# ═══════════════════════════════════════════════════════════
# Slide 14: 横展、断点佐证
# ═══════════════════════════════════════════════════════════
slide14 = prs.slides[13]
for shape in slide14.shapes:
    if shape.has_text_frame:
        txt = shape.text.strip()
        if "横展了哪些项目" in txt or "横展项目" in txt:
            clear_textbox(shape)
            tf = shape.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "横展项目："
            run.font.bold = True
            for item in ["全规格铜箔产品", "所有包装班组/工位", "所有客户端出货标准统一"]:
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = item
                run.font.size = Pt(10)
        elif "长期措施断点批次" in txt or "断点批次" in txt:
            clear_textbox(shape)
            tf = shape.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "断点批次："
            run.font.bold = True
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = "2026-05-31及以后生产/包装批次"
            run.font.size = Pt(10)
            run.font.bold = True
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = "（改善措施全面落地，无同类人为端面划伤风险）"
            run.font.size = Pt(9)

print('Slide 14 done')

# ═══════════════════════════════════════════════════════════
# 保存
# ═══════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f'\nAll done! Saved to: {OUTPUT}')
